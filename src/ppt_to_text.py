#!/usr/bin/env python
# Install dependency: pip install python-pptx

from pptx import Presentation
import glob

FILE_NAME = "example.pdf"

for eachfile in glob.glob():
    prs = Presentation(eachfile)
    print(eachfile)
    print("----------------------")
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)
